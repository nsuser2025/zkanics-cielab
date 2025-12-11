import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

def mkpptx_gui(df, images, result):
    template_form = st.radio("テンプレートを使用しますか？", ["Yes", "No"], index=1, horizontal=True)
    
    # テンプレートの読み込み
    if template_form == "Yes":
        uploaded_template = st.file_uploader("PPTXテンプレートをアップロード", type=["pptx"])
        if uploaded_template is None:
            st.warning("PPTXテンプレートをアップロードしてください。")
            return
        template_bytes = BytesIO(uploaded_template.read())
        prs = Presentation(template_bytes)
    else:
        # テンプレート読み込みのフォールバックは健全だが、
        # GitHub上のファイルを BytesIO 経由で読み込む方がより堅牢 (外部リソース読み込みコードが必要)
        try:
            # 現在の動作を尊重しつつ、エラー時の組み込みテンプレートへフォールバック
            prs = Presentation("./default_template.pptx")
        except FileNotFoundError:
            prs = Presentation() # 組み込みの標準テンプレートを使用
        
    # レイアウト設定
    rows, cols_num = 2, 3
    width = Inches(2.8) 
    height = Inches(2.1)
    spacing_x, spacing_y = Inches(0.2), Inches(0.6)
    left_margin, bottom_margin = Inches(0.7), Inches(0.8)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    max_per_slide = rows * cols_num

    for batch_start in range(0, len(result), max_per_slide):
        batch = result[batch_start: batch_start + max_per_slide]
        
        slide_layout_index = min(2, len(prs.slide_layouts) - 1)
        slide_layout = prs.slide_layouts[slide_layout_index]
        image_slide = prs.slides.add_slide(slide_layout)

        total_height = rows * height + (rows - 1) * spacing_y

        for idx, name in enumerate(batch):
            if name in images:
                image = images[name]
                
                # BytesIO 改良点 1: PIL Image をメモリバッファ (BytesIO) に保存
                image_stream = BytesIO()
                # add_pictureが識別しやすいようにPNG形式で保存することが多い
                image.save(image_stream, format='PNG') 
                image_stream.seek(0) # ストリームのポインタを先頭に戻す

                row = idx // cols_num
                col = idx % cols_num
                left = left_margin + col * (width + spacing_x)
                top = slide_height - total_height - bottom_margin + row * (height + spacing_y)
                
                # BytesIO 改良点 2: add_picture に BytesIO オブジェクトを直接渡す
                try:
                    image_slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
                except Exception as e:
                    st.warning(f"画像 '{name}' の追加に失敗しました。BytesIO形式の問題か、ファイル形式が正しくありません。エラー: {e}")
                    continue

                # ... (条件テキスト作成ロジック) ...
                condition_row = df[df["ファイル名"] == name]
                if not condition_row.empty:
                    condition_row = condition_row.iloc[0]
                    cond_text = (f"{condition_row['試験']}-"
                                 f"{condition_row['測定面']}-"
                                 f"{condition_row['正極']}-"
                                 f"{condition_row['測定']}-"
                                 f"{condition_row['電解液']}-"
                                 f"{condition_row['倍率']}") 

                    textbox_left = left
                    textbox_top = top + height + Inches(0.05) 
                    textbox_width = width
                    # Inches 統一: テキストボックスの高さを Inches で指定
                    textbox_height = Inches(0.5) 
                    textbox = image_slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = True
                    text_frame.clear()
                    
                    lines = cond_text.split("\n")
                    if lines:
                        p = text_frame.paragraphs[0]
                        p.text = lines[0]
                        p.font.size = Pt(9)
                        for line in lines[1:]:
                            p = text_frame.add_paragraph()
                            p.text = line
                            p.font.size = Pt(9)

    # BytesIO 改良点 3: 最終PPTXファイルをメモリ (BytesIO) に保存
    output = BytesIO()
    try:
        prs.save(output)
    except Exception as e:
        st.error(f"PPTXファイルの保存中にエラーが発生しました: {e}")
        return

    output.seek(0)
    st.success("PPTXファイルの準備ができました。")
    
    # クリーンアップ不要: ディスクI/Oがないため、残ファイル削除コードは不要

    # ダウンロード
    st.download_button(
        "PPTXをダウンロード",
        data=output,
        file_name="output.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
